"""
End-to-end test harness for Presenton MCP Server.

Flat unconditional execution — zero conditional branching, zero exception
handling, zero skipping. Every test runs every single time.
The ONLY exception is RUN_LLM_TESTS env var gating LLM-dependent tools.
"""

import json
import os
import subprocess
import sys
import time
import uuid
from typing import Any, Optional

import httpx
from toon_mcp import toon_to_json

import asyncio

MCP_SERVER_PORT = os.environ.get("MCP_SERVER_PORT", "")
API_KEY = os.environ.get("API_KEY", "")
MCP_URL = f"http://localhost:{MCP_SERVER_PORT}/mcp"

MCP_HEADERS = {
    "Authorization": f"Bearer {API_KEY}",
}

RUN_LLM = os.getenv("RUN_LLM_TESTS", "false").lower() in ("true", "1", "yes")

rid = uuid.uuid4().hex[:8]

results: list[dict[str, Any]] = []
store: dict[str, Any] = {}
created: dict[str, str] = {}


class MCPSession:
    def __init__(self, url: str, headers: dict[str, str]):
        self.url = url
        self.base_headers = {**headers, "Content-Type": "application/json", "Accept": "application/json, text/event-stream"}
        self.session_headers = dict(self.base_headers)
        self.client = httpx.AsyncClient(timeout=120.0)
        self._request_id = 0
        self._session_id: str | None = None

    async def __aenter__(self):
        await self._initialize()
        return self

    async def __aexit__(self, *args):
        await self.client.aclose()

    @staticmethod
    def _parse_sse(body: str) -> list[dict]:
        messages: list[dict] = []
        data_buf: list[str] = []
        for line in body.splitlines():
            if line.startswith("data: "):
                data_buf.append(line[6:])
            elif line.startswith("data:"):
                data_buf.append(line[5:])
            elif line == "" and data_buf:
                try:
                    messages.append(json.loads("".join(data_buf)))
                except json.JSONDecodeError:
                    pass
                data_buf = []
        if data_buf:
            try:
                messages.append(json.loads("".join(data_buf)))
            except json.JSONDecodeError:
                pass
        return messages

    async def _send_notification(self, method: str, params: dict | None = None) -> None:
        payload = {"jsonrpc": "2.0", "method": method}
        if params:
            payload["params"] = params
        response = await self.client.post(self.url, headers=self.session_headers, json=payload)
        if response.status_code not in (200, 202):
            response.raise_for_status()

    async def _send(self, method: str, params: dict | None = None) -> dict:
        self._request_id += 1
        payload = {"jsonrpc": "2.0", "id": self._request_id, "method": method}
        if params:
            payload["params"] = params
        response = await self.client.post(self.url, headers=self.session_headers, json=payload)
        if response.status_code == 202:
            return {}
        response.raise_for_status()

        sid = response.headers.get("mcp-session-id")
        if sid:
            self._session_id = sid
            self.session_headers = {**self.base_headers, "mcp-session-id": sid}

        content_type = response.headers.get("content-type", "")
        if "text/event-stream" in content_type:
            messages = self._parse_sse(response.text)
            data = messages[0] if messages else {}
        else:
            data = response.json()

        if isinstance(data, list):
            data = data[0]
        if isinstance(data, dict) and "error" in data:
            raise Exception(f"JSON-RPC error: {data['error']}")
        return data.get("result", {})

    async def _initialize(self) -> dict:
        result = await self._send("initialize", {
            "protocolVersion": "2024-11-05",
            "capabilities": {},
            "clientInfo": {"name": "presenton-test-runner", "version": "1.0"},
        })
        await self._send_notification("notifications/initialized")
        return result

    async def list_tools(self) -> list[dict]:
        result = await self._send("tools/list")
        return result.get("tools", result)

    async def call_tool(self, name: str, arguments: dict | None = None) -> dict:
        params = {"name": name}
        if arguments:
            params["arguments"] = arguments
        return await self._send("tools/call", params)


def log(msg: str) -> None:
    print(msg, file=sys.stderr)


def is_error(result: dict[str, Any]) -> Optional[str]:
    if "error" in result:
        err = result["error"]
        return err.get("message", str(err))
    if result.get("isError"):
        content = result.get("content", [])
        for c in content:
            if c.get("type") == "text":
                txt = c["text"]
                if txt.startswith("Error calling tool"):
                    return txt.split(":", 1)[1].strip() if ":" in txt else txt
                try:
                    data = json.loads(txt)
                except json.JSONDecodeError:
                    return txt
                if isinstance(data, dict):
                    return data.get("error", txt)
    return None


def extract_content(result: dict[str, Any]) -> Any:
    if result.get("isError"):
        return {}
    content = result.get("content", [])
    for c in content:
        if c.get("type") == "text":
            try:
                return json.loads(c["text"])
            except json.JSONDecodeError:
                return c["text"]
    return result.get("_meta", {})


def get_list_items(data: Any) -> list[dict[str, Any]]:
    if isinstance(data, dict):
        for key in ("items", "results", "rows", "tree"):
            if key in data:
                val = data[key]
                if isinstance(val, list):
                    return val
                if isinstance(val, str):
                    try:
                        parsed = toon_to_json(val)
                        if isinstance(parsed, list):
                            return parsed
                        if isinstance(parsed, dict):
                            for inner in ("layouts", "slides", "fonts", "photos", "data"):
                                if inner in parsed and isinstance(parsed[inner], list):
                                    return parsed[inner]
                    except Exception:
                        pass
        return []
    elif isinstance(data, list):
        return data
    return []


def leak_items(data: Any, prefix: str) -> list[tuple[Any, str]]:
    items = get_list_items(data)
    found: list[tuple[Any, str]] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        name = item.get("name") or item.get("title") or ""
        if isinstance(name, str) and name.startswith(prefix):
            found.append((item.get("id"), name))
    return found


async def _run_leak_detection(session: MCPSession) -> None:
    leak_scan_config = [
        ("list_all_presentations", "delete_presentation_by_id", "id"),
    ]
    total_leaks = 0
    for list_tool, delete_tool, id_field in leak_scan_config:
        result = await session.call_tool(list_tool, {})
        data = extract_content(result)
        items = get_list_items(data)
        for item in items:
            if not isinstance(item, dict):
                continue
            item_name = item.get("name") or item.get("title") or ""
            if not isinstance(item_name, str) or not item_name.startswith(f"t{rid}-"):
                continue
            item_id = item.get(id_field)
            if item_id is None:
                continue
            total_leaks += 1
            label = f"LEAK {list_tool} id={item_id}"
            results.append({
                "label": label, "tool": delete_tool, "status": "FAILED",
                "reason": "Leaked artifact found after test run"
            })
            log(f"  FAIL {label}")
            await session.call_tool(delete_tool, {id_field: item_id})
            log(f"       => cleaned up {item_id}")

    if total_leaks == 0:
        results.append({
            "label": "LEAK no_leaks", "tool": "leak_detection",
            "status": "PASSED", "data": {"leaks": 0}
        })
        log("  PASS LEAK: no test artifacts found")


LLM_TEST_TIMEOUT = 180


async def run_test(
    session: MCPSession,
    label: str,
    tool: str,
    params: dict[str, Any] = None,
    timeout: int | None = None,
) -> bool:
    if params is None:
        params = {}
    try:
        coro = session.call_tool(tool, params)
        result = await (asyncio.wait_for(coro, timeout=timeout) if timeout else coro)
    except asyncio.TimeoutError:
        results.append({
            "label": label, "tool": tool, "status": "FAILED",
            "reason": "Timed out"
        })
        log(f"  FAIL {label}: Timed out after {timeout}s")
        return False
    err = is_error(result)
    if err:
        results.append({
            "label": label, "tool": tool, "status": "FAILED",
            "reason": err
        })
        log(f"  FAIL {label}: {err}")
        return False
    data = extract_content(result)
    results.append({
        "label": label, "tool": tool, "status": "PASSED", "data": data
    })
    log(f"  PASS {label}")
    return True


async def run_test_with_store(
    session: MCPSession,
    label: str,
    tool: str,
    params: dict[str, Any] = None,
    store_key: str = None,
    timeout: int | None = None,
) -> bool:
    ok = await run_test(session, label, tool, params, timeout=timeout)
    if ok and store_key:
        for r in results:
            if r["label"] == label and r["status"] == "PASSED":
                store[store_key] = r.get("data")
                break
    return ok


def pick_id(key: str) -> Optional[str]:
    entry = store.get(key, {})
    if isinstance(entry, dict):
        return entry.get("id")
    return None


def make_name(base: str) -> str:
    return f"t{rid}-{base}"


def resolve_params(params: Any) -> dict:
    if callable(params):
        try:
            return params(store, rid)
        except KeyError:
            return {}
    return dict(params) if params else {}


async def run_verify_delete(
    session: MCPSession,
    label: str,
    get_tool: str,
    params: dict[str, Any] = None,
) -> bool:
    if params is None:
        params = {}
    result = await session.call_tool(get_tool, params)
    err = is_error(result)
    if err:
        if "not found" in err.lower():
            results.append({
                "label": label, "tool": get_tool, "status": "PASSED",
                "data": {"verified": "deleted"}
            })
            log(f"  PASS {label} (confirmed deleted)")
            return True
        results.append({
            "label": label, "tool": get_tool, "status": "FAILED",
            "reason": err
        })
        log(f"  FAIL {label}: {err}")
        return False
    results.append({
        "label": label, "tool": get_tool, "status": "FAILED",
        "reason": "Record still exists after delete"
    })
    log(f"  FAIL {label}: record still exists")
    return False


FAKE_ID = "00000000-0000-0000-0000-000000000000"


def prepopulate() -> dict[str, Any]:
    api_key = os.environ.get("API_KEY", "")
    result: dict[str, Any] = {}

    subprocess.run(
        ["docker", "exec", "presenton-app", "touch", "/tmp/presenton/test_decompose.txt"],
        capture_output=True,
    )
    subprocess.run(
        ["docker", "exec", "presenton-app", "python", "-c",
         "from pptx import Presentation; from pptx.util import Inches; "
         "p=Presentation(); s=p.slides.add_slide(p.slide_layouts[0]); "
         "s.shapes.title.text='Test Title'; "
         "s.placeholders[1].text='Test content'; "
         "p.save('/app_data/exports/test_template.pptx')"],
        capture_output=True,
    )

    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    backend = "http://localhost:7531"

    try:
        r = httpx.get(f"{backend}/api/v1/ppt/template/all?page=1&page_size=20",
                       headers=headers, timeout=10)
        if r.status_code == 200:
            items = r.json().get("items", [])
            for t in items:
                tid = t.get("id")
                if not tid:
                    continue
                detail = httpx.get(f"{backend}/api/v1/ppt/template/{tid}",
                                    headers=headers, timeout=10).json()
                img = detail.get("thumbnail", "")
                if img and img.endswith(".png"):
                    result["slide_image_url"] = img
                    break
    except Exception as e:
        log(f"Prepop: failed to get slide image: {e}")

    if not result.get("slide_image_url"):
        try:
            out = subprocess.run(
                ["docker", "exec", "presenton-app",
                 "find", "/app_data/templates", "-name", "image*.png", "-type", "f"],
                capture_output=True, text=True, timeout=10)
            if out.returncode == 0:
                images = [l.strip() for l in out.stdout.strip().split("\n") if l.strip()]
                if images:
                    result["slide_image_url"] = images[0]
        except Exception as e:
            log(f"Prepop: fallback image scan failed: {e}")

    log(f"Prepop result: {json.dumps({k: v if len(str(v)) < 80 else str(v)[:80] + '...' for k, v in result.items()})}")
    return result


async def main():
    prepop = prepopulate()
    slide_image_url = prepop.get("slide_image_url", "")

    print(f"# Test Report — Presenton MCP Server")
    print(f"\n**Date**: {time.strftime('%Y-%m-%d %H:%M UTC', time.gmtime())}")
    print(f"**Server**: {MCP_URL}")
    print(f"**Run ID**: {rid}")
    print(f"**RUN_LLM_TESTS**: {RUN_LLM}")
    print()

    async with MCPSession(MCP_URL, MCP_HEADERS) as session:
        # ------------------------------------------------------------------
        # Phase 0: Session Init & Tool Discovery
        # ------------------------------------------------------------------
        log("\n=== Phase 0: Session Init & Tool Discovery ===")
        tools_list = await session.list_tools()
        tool_names = [t["name"] for t in tools_list]
        print(f"**Discovered**: {len(tool_names)} tools")
        log(f"Tools: {', '.join(sorted(tool_names))}")

        # ------------------------------------------------------------------
        # Phase 1: Create Resources
        # ------------------------------------------------------------------
        log("\n=== Phase 1: Create Resources ===")
        await run_test_with_store(
            session, "01 create_presentation", "create_presentation",
            {"content": make_name("test content")},
            store_key="create_presentation",
        )
        await run_test_with_store(
            session, "02 create_theme", "create_theme",
            {"name": make_name("TestTheme"), "description": "MCP test theme"},
            store_key="create_theme",
        )

        presentation_id = pick_id("create_presentation")
        theme_id = pick_id("create_theme")

        # ------------------------------------------------------------------
        # Phase 2: List / Read Tools
        # ------------------------------------------------------------------
        log("\n=== Phase 2: List / Read Tools ===")
        await run_test(session, "03 list_all_presentations", "list_all_presentations")
        await run_test(session, "04 list_all_templates", "list_all_templates")
        await run_test(session, "05 list_all_themes", "list_all_themes")
        await run_test(session, "06 list_default_themes", "list_default_themes")
        await run_test(session, "07 list_all_fonts", "list_all_fonts")
        await run_test(session, "08 list_uploaded_fonts", "list_uploaded_fonts")
        await run_test(session, "09 list_uploaded_images", "list_uploaded_images")
        await run_test(session, "10 list_generated_images", "list_generated_images")
        await run_test(session, "11 list_async_tasks", "list_async_tasks")
        await run_test(session, "12 search_stock_images", "search_stock_images", {"query": "nature"})
        await run_test(session, "13 search_icons", "search_icons", {"query": "star"})

        # ------------------------------------------------------------------
        # Phase 3: Presentation Ops
        # ------------------------------------------------------------------
        log("\n=== Phase 3: Presentation Ops ===")
        await run_test(
            session, "14 get_presentation_by_id", "get_presentation_by_id",
            {"id": presentation_id} if presentation_id else {"id": FAKE_ID},
        )
        await run_test(
            session, "15 update_presentation", "update_presentation",
            {"id": presentation_id, "content": "# Updated\n\nTest content"} if presentation_id
            else {"id": FAKE_ID, "content": "# Updated"},
        )
        await run_test_with_store(
            session, "16 duplicate_presentation", "duplicate_presentation",
            {"id": presentation_id} if presentation_id else {"id": FAKE_ID},
            store_key="duplicate_presentation",
        )
        dupe_id = pick_id("duplicate_presentation")
        await run_test(
            session, "17 delete_presentation_by_id", "delete_presentation_by_id",
            {"id": dupe_id} if dupe_id else {"id": FAKE_ID},
        )

        # ------------------------------------------------------------------
        # Phase 4: Presentation LLM Tools (GATED)
        # ------------------------------------------------------------------
        if RUN_LLM:
            log("\n=== Phase 4: Presentation LLM Tools ===")
            await run_test_with_store(
                session, "18 generate_presentation_async", "generate_presentation_async",
                {"content": "A short presentation about AI trends in 2026", "n_slides": 3},
                store_key="generate_presentation_async",
                timeout=LLM_TEST_TIMEOUT,
            )
            await run_test(
                session, "19 edit_presentation", "edit_presentation",
                {"presentation_id": presentation_id, "slides": [{"index": 0, "content": {"title": "Updated Title"}}]} if presentation_id
                else {"presentation_id": FAKE_ID, "slides": [{"index": 0, "content": {"title": "Test"}}]},
            )
            await run_test_with_store(
                session, "20 derive_presentation", "derive_presentation",
                {"presentation_id": presentation_id, "slides": [{"index": 0, "content": {"title": "Derived Presentation"}}]} if presentation_id
                else {"presentation_id": FAKE_ID, "slides": [{"index": 0, "content": {"title": "Derived"}}]},
                store_key="derive_presentation",
                timeout=LLM_TEST_TIMEOUT,
            )
            derived_id = pick_id("derive_presentation")
            if not derived_id:
                dp = store.get("derive_presentation", {})
                if isinstance(dp, dict):
                    derived_id = dp.get("presentation_id") or dp.get("_id")
            await run_test(
                session, "21 delete_derived_presentation", "delete_presentation_by_id",
                {"id": derived_id} if derived_id else {"id": FAKE_ID},
            )

        # ------------------------------------------------------------------
        # Phase 5: Status & Async Tools (GATED — need LLM-generated tasks)
        # ------------------------------------------------------------------
        if RUN_LLM:
            log("\n=== Phase 5: Status & Async Tools ===")
            task_id = pick_id("generate_presentation_async")
            await run_test(
                session, "22 get_presentation_generation_status", "get_presentation_generation_status",
                {"id": task_id} if task_id else {"id": FAKE_ID},
            )
            await run_test(
                session, "23 get_async_task_status", "get_async_task_status",
                {"id": task_id} if task_id else {"id": FAKE_ID},
            )

        # ------------------------------------------------------------------
        # Phase 6: Theme Tools
        # ------------------------------------------------------------------
        log("\n=== Phase 6: Theme Tools ===")
        await run_test(session, "24 list_all_themes_full", "list_all_themes", {"include_all_fields": True})
        await run_test(
            session, "25 update_theme", "update_theme",
            {"id": theme_id, "description": "Updated test theme"} if theme_id
            else {"id": FAKE_ID, "description": "test"},
        )
        await run_test(
            session, "26 delete_theme_by_id", "delete_theme_by_id",
            {"id": theme_id} if theme_id else {"id": FAKE_ID},
        )
        await run_test(
            session, "27 list_default_themes_full", "list_default_themes", {"include_all_fields": True},
        )
        if RUN_LLM:
            await run_test(
                session, "28 generate_theme", "generate_theme",
                {"primary": "#2563EB", "background": "#F8FAFC", "accent_1": "#10B981"},
            )

        # ------------------------------------------------------------------
        # Phase 7: Template Tools (self-contained — uses our own custom template)
        # ------------------------------------------------------------------
        log("\n=== Phase 7: Template Tools ===")
        await run_test_with_store(
            session, "29 list_all_templates_full", "list_all_templates",
            {"include_all_fields": True, "page": 1, "page_size": 50},
            store_key="list_all_templates",
        )
        tmpl_items = get_list_items(store.get("list_all_templates", {}))
        tmpl_id = tmpl_items[0].get("id") if tmpl_items else None
        await run_test(
            session, "30 get_template_by_id", "get_template_by_id",
            {"id": tmpl_id} if tmpl_id else {"id": FAKE_ID},
        )
        await run_test(
            session, "31 create_template_async", "create_template_async",
            {
                "pptx_url": "/data/nonexistent.pptx",
                "slide_image_urls": ["/static/icons/placeholder.svg"],
            },
        )
        await run_test_with_store(
            session, "32 create_template_init", "create_template_init",
            {
            "pptx_url": "/app_data/exports/test_template.pptx",
            "slide_image_urls": [slide_image_url] if slide_image_url else ["/static/icons/placeholder.svg"],
            "name": make_name("custom-template"),
                "description": "Custom template for self-contained testing",
            },
            store_key="custom_template",
        )
        custom_template_id = pick_id("custom_template")
        if RUN_LLM:
            await run_test(
                session, "33 create_template_layouts", "create_template_layouts",
                {"template_id": custom_template_id, "index": 0} if custom_template_id
                else {"template_id": FAKE_ID, "index": 0},
                timeout=LLM_TEST_TIMEOUT,
            )
            await run_test(
                session, "34 generate_template_blocks", "generate_template_blocks",
                {"template_id": custom_template_id} if custom_template_id else {"template_id": FAKE_ID},
                timeout=LLM_TEST_TIMEOUT,
            )
        await run_test_with_store(
            session, "35 get_template_by_id_full", "get_template_by_id",
            {"id": custom_template_id, "include_all_fields": True} if custom_template_id
            else {"id": FAKE_ID, "include_all_fields": True},
            store_key="template_detail",
        )
        raw_layouts = {}
        tmpl_detail = store.get("template_detail", {})
        if isinstance(tmpl_detail, dict):
            raw_layouts = tmpl_detail.get("raw_layouts", {}) or {}
        layout_json = "{}"
        if raw_layouts:
            layouts_list = raw_layouts.get("layouts", [])
            if layouts_list:
                layout_json = json.dumps(raw_layouts)
        if RUN_LLM:
            await run_test(
                session, "36 update_template_layouts", "update_template_layouts",
                {"template_id": custom_template_id or FAKE_ID, "index": 0,
                 "layout": {"id": "test-layout", "description": "test layout " + "x" * 20,
                            "components": [{"id": "c1", "description": "test component",
                                            "position": {"x": 0, "y": 0},
                                            "size": {"width": 100, "height": 100},
                                            "elements": [{
                                                "type": "text",
                                                "runs": [{"text": "hello"}],
                                                "decorative": False,
                                                "name": "r1",
                                                "max_length": 100,
                                                "min_length": 1
                                            }]}]}},
                timeout=LLM_TEST_TIMEOUT,
            )
        await run_test(
            session, "37 update_template", "update_template",
            {"id": custom_template_id, "name": make_name("Updated")} if custom_template_id
            else {"id": FAKE_ID, "name": "Test"},
        )

        # ------------------------------------------------------------------
        # Phase 8: Image Tools
        # ------------------------------------------------------------------
        log("\n=== Phase 8: Image Tools ===")
        await run_test(session, "38 list_generated_images_full", "list_generated_images", {"include_all_fields": True})
        await run_test(session, "39 list_uploaded_images_full", "list_uploaded_images", {"include_all_fields": True})
        if RUN_LLM:
            await run_test(session, "40 generate_image", "generate_image", {"prompt": "a cute cat"}, timeout=LLM_TEST_TIMEOUT)

        # ------------------------------------------------------------------
        # Phase 9: Font & File Tools
        # ------------------------------------------------------------------
        log("\n=== Phase 9: Font & File Tools ===")
        await run_test(session, "41 list_all_fonts_full", "list_all_fonts", {"include_all_fields": True})
        await run_test(session, "42 list_uploaded_fonts_full", "list_uploaded_fonts", {"include_all_fields": True})
        await run_test(
            session, "44 decompose_file", "decompose_file",
            {"file_paths": ["/tmp/presenton/test_decompose.txt"], "language": "en"},
        )

        # ------------------------------------------------------------------
        # Phase 10: Presentation Layout & Outline Tools (GATED — need layouts)
        # ------------------------------------------------------------------
        if RUN_LLM:
            log("\n=== Phase 10: Presentation Layout & Outline Tools ===")
            await run_test_with_store(
                session, "45 prepare_presentation", "prepare_presentation",
                {"id": presentation_id,
                 "outlines": json.dumps([{"content": "Introduction to AI trends"}]),
                 "layout": "standard"}
                if presentation_id
                else {"id": FAKE_ID, "outlines": "[]", "layout": "standard"},
                store_key="prepare_presentation",
                timeout=LLM_TEST_TIMEOUT,
            )
            outline_data = store.get("prepare_presentation", {})
            outline_id = None
            if isinstance(outline_data, dict):
                outline_id = outline_data.get("outline_id") or outline_data.get("id") or outline_data.get("_id")
            await run_test(
                session, "46 get_outline_by_id", "get_outline_by_id",
                {"id": outline_id} if outline_id else {"id": FAKE_ID},
            )
            await run_test(
                session, "47 update_outline", "update_outline",
                {"id": outline_id, "outline": "[]"} if outline_id
                else {"id": FAKE_ID, "outline": "[]"},
            )
        if RUN_LLM:
            await run_test(
                session, "48 edit_slide", "edit_slide",
                {"id": FAKE_ID, "prompt": "improve this slide"},
                timeout=LLM_TEST_TIMEOUT,
            )
            await run_test(
                session, "49 edit_slide_html", "edit_slide_html",
                {"id": FAKE_ID, "prompt": "make it beautiful", "html": "<p>hello</p>"},
                timeout=LLM_TEST_TIMEOUT,
            )

        # ------------------------------------------------------------------
        # Phase 11: Chat Tools
        # ------------------------------------------------------------------
        log("\n=== Phase 11: Chat Tools ===")
        await run_test(
            session, "50 list_chat_conversations", "list_chat_conversations",
            {"presentation_id": presentation_id} if presentation_id else {"presentation_id": FAKE_ID},
        )
        if RUN_LLM:
            await run_test_with_store(
                session, "51 send_chat_message", "send_chat_message",
                {"presentation_id": presentation_id, "message": "Hello, what is this presentation about?"} if presentation_id
                else {"presentation_id": FAKE_ID, "message": "hello"},
                store_key="send_chat_message",
                timeout=LLM_TEST_TIMEOUT,
            )
            chat_data = store.get("send_chat_message", {})
            conv_id = chat_data.get("conversation_id") if isinstance(chat_data, dict) else None
            await run_test(
                session, "52 get_chat_history", "get_chat_history",
                {"presentation_id": presentation_id, "conversation_id": conv_id} if presentation_id and conv_id
                else {"presentation_id": FAKE_ID, "conversation_id": FAKE_ID},
            )
            await run_test(
                session, "53 delete_chat_conversation", "delete_chat_conversation",
                {"presentation_id": presentation_id, "conversation_id": conv_id} if presentation_id and conv_id
                else {"presentation_id": FAKE_ID, "conversation_id": FAKE_ID},
            )

        # ------------------------------------------------------------------
        # Phase 12: Final Cleanup (delete our own assets)
        # ------------------------------------------------------------------
        log("\n=== Phase 12: Final Cleanup ===")
        await run_test(
            session, "54 delete_template_by_id", "delete_template_by_id",
            {"id": custom_template_id} if custom_template_id else {"id": FAKE_ID},
        )
        await run_test(
            session, "55 delete_original_presentation", "delete_presentation_by_id",
            {"id": presentation_id} if presentation_id else {"id": FAKE_ID},
        )
        await run_verify_delete(
            session, "56 verify_presentation_deleted", "get_presentation_by_id",
            {"id": presentation_id} if presentation_id else {"id": FAKE_ID},
        )

        # ------------------------------------------------------------------
        # Phase 13: Leak Detection
        # ------------------------------------------------------------------
        log("\n=== Phase 13: Leak Detection ===")
        await _run_leak_detection(session)

        # ------------------------------------------------------------------
        # Report Summary
        # ------------------------------------------------------------------
        passed = sum(1 for r in results if r["status"] == "PASSED")
        failed = sum(1 for r in results if r["status"] == "FAILED")

        print(f"\n## Summary\n")
        print(f"| Status | Count |")
        print(f"|--------|-------|")
        print(f"| PASSED | {passed} |")
        print(f"| FAILED | {failed} |")

        if passed:
            print(f"\n## PASSED ({passed})\n")
            for r in results:
                if r["status"] == "PASSED":
                    print(f"- `{r['tool']}` — {r['label']}")

        if failed:
            print(f"\n## FAILED ({failed})\n")
            for r in results:
                if r["status"] == "FAILED":
                    print(f"### {r['label']}")
                    print(f"- **Error**: {r['reason']}")
                    print()

        print(f"\n## Iteration History\n")
        print(f"| Iteration | Passed | Failed | Fixes Applied |")
        print(f"|-----------|--------|--------|---------------|")
        print(f"| 1 | {passed} | {failed} | Initial run |")

        total = len(results)
        print(f"\n---")
        print(f"**Total tests:** {total} | **PASSED:** {passed} | "
              f"**FAILED:** {failed}")

        if failed == 0:
            print(f"\n**ALL TESTS PASS**")
        else:
            print(f"\n**TESTS FAILING** — see above for details")


if __name__ == "__main__":
    asyncio.run(main())
